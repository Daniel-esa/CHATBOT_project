import pandas as pd
from kor.extraction import create_extraction_chain
from kor.nodes import Object, Text, Number
from PyPDF2 import PdfReader
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import io





def load_pdf(file):
    """
    Extracts text from a PDF document.

    :param file_path: str, path to the PDF file
    :return: str, extracted text
    """
    # Create a PDF reader object
    reader = PdfReader(file)
   
    # Extract text
    extracted_text = []
   
    for page in reader.pages:
        extracted_text.append(page.extract_text())
   
    return "\n".join(extracted_text)

def load_pdf_minder(file):
    
    text = extract_text(file)
    return text


def load_ppt(file):
    """
    Extracts text from a PowerPoint presentation.

    :param file_path: str, path to the PowerPoint file
    :return: str, extracted text
    """

    # Load the presentation
    presentation = Presentation(io.BytesIO(file.read()))
   
    # Extract text
    extracted_text = []
   
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text.append(shape.text)
   
    return "\n".join(extracted_text)

def load_docx(file):
    """
    Extracts text from a Word document (.docx).

    :param file_path: str, path to the Word document
    :return: str, extracted text
    """

    # Load the document
    document = Document(io.BytesIO(file.read()))
   
    # Extract text
    extracted_text = []
   
    for para in document.paragraphs:
        extracted_text.append(para.text)
   
    return "\n".join(extracted_text)

def load_txt(file):
    """
    Extracts text from a text file (.txt).

    :param file_path: str, path to the text file
    :return: str, extracted text
    """    
    stringio = StringIO(file.getvalue().decode("utf-8"))
    # To read file as string:
    extracted_text = stringio.read()
    return extracted_text


def load_document(file):
    
    if file.name.endswith('.pdf'):
        text = load_pdf_minder(file)
    elif file.name.endswith('.docx'):
        text = load_docx(file)
    elif file.name.endswith('.txt'):
        text = load_txt(file)
    elif file.name.endswith('.pptx'):
        text = load_ppt(file)
    return text

def printOutput(output):
    print(json.dumps(output,sort_keys=True, indent=3))



def extract_fields(names, descriptions, exs1, exs2, exs3, llm, source_text, language="English"):


    #16 elt 0 -> 15

    example_1 = {
        "Supplier":"IBM",
        "Start date":"2020-02-20",
        "End date":"2024 Aug 10 ",
        "Renewal Types":"not renewable",
        "Renewal Period":"3 months",
        "Contract global amount":"The total contract amount is set at €500,000 (five hundred thousand euros) and shall be paid in euros throughout the duration of the contract1.",
        "Contract currency":"€",
        "National Regulation / governing law":"French Law",
        "Competent jurisdiction":"Tribunal de Nanterre",
        "Personal data law":"California Consumer Privacy Act",
        "GDPR":"General Data Protection Regulation",
        "SRB":"Single Resolution Board",
        "DORA":"Digital Operational Resilience Act",
        "EBA":"European Banking Authority",
        "Reversibility":"In the event of contract termination, the Supplier agrees to provide all necessary assistance to transfer the services to a new provider, ensuring minimal disruption to the Client's operations",
        "Substituability":"The Supplier agrees to provide all necessary support and documentation to facilitate the substitution of services by another provider, ensuring a smooth transition and compliance with contractual obligations",
    }
    
 
    
    

    schema = Object(
        id="contract",
        description="Information about a contract",
        attributes = [
            Text(
                id="Supplier",
                description="It can be an external supplier or an internal provider in case of intra-group contracts. You can usually find this data at the beginning of a contract ('the parties').",
            ),
            Text(
                id="Start date",
                description="Contract start date.",
            ),
            Text(
                id="End date",
                description="Contract end date.",
            ),
            Text(
                id="Renewal Types",
                description="3 options possible: - not renewable; - renewable automatically under the same conditions (tacit reconduction);- renewable with negociation of the terms",
            ),
            Text(
                id="Renewal Period",
                description="Period of contract extension upon renewal (usually in months but not always)",
            ),
            Text(
                id="Contract global amount",
                description="Overall amount of the commitment indicated in the contract over its entire period of validity",
            ),
            Text(
                id="Contract currency",
                description="Currency of the contract global amount",
            ),
            Text(
                id="National Regulation / governing law",
                description="National regulation / governing law : the legal system under which the contract will be interpreted. Not only national (can also be states regulation) / (droit applicable en FR)",
            ),
            Text(
                id="Competent jurisdiction",
                description="Competent jurisdiction : the specific court that have the authority to resolve disputes (for example: Tribunal de Nanterre)",
            ),
            Text(
                id="Personal data law",
                description="Sentences related to a personal data law",
            ),
            Text(
                id="GDPR",
                description="Sentences related  to GDPR",
            ),
            Text(
                id="SRB",
                description="Sentences related SRB or Single Resolution Board. Key words : SRB, Single Resolution Board, resolution, resolution plan, service continuity, regulatory…",
            ),
            Text(
                id="DORA",
                description="Sentences related to DORA or Digital Operational Resilience Act. Key words : DORA,- Digital Operational Resilience Act, operational resilience, incident, cooperation,  information sharing, regulatory, authorities, remedial actions…",
            ),
            Text(
                id="EBA",
                description="Sentences related to EBA or European Banking Authority. Key words : EBA, European Banking Authority, EBA guidelines",
            ),
            Text(
                id="Reversibility",
                description="Sentences related to contract reversibility.",
            ),
            Text(
                id="Substituability",
                description="Sentences related to substituability.",
            )
        ],
        examples=[
            (
                source_text,
                [
                    example_1,
                ],
            )
        ]
    )



    chain = create_extraction_chain(llm, schema)
    result = chain.invoke(source_text)
    return result["data"]
    #return example_1
