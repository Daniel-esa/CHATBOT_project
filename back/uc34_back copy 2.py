import pandas as pd
from kor.extraction import create_extraction_chain
from kor.nodes import Object, Text, Number
from PyPDF2 import PdfReader
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import io





def load_pdf(file):
    """
    Extracts text from a PDF document.

    :param file_path: str, path to the PDF file
    :return: str, extracted text
    """
    # Create a PDF reader object
    reader = PdfReader(file)
   
    # Extract text
    extracted_text = []
   
    for page in reader.pages:
        extracted_text.append(page.extract_text())
   
    return "\n".join(extracted_text)

def load_pdf_minder(file):
    
    text = extract_text(file)
    return text


def load_ppt(file):
    """
    Extracts text from a PowerPoint presentation.

    :param file_path: str, path to the PowerPoint file
    :return: str, extracted text
    """

    # Load the presentation
    presentation = Presentation(io.BytesIO(file.read()))
   
    # Extract text
    extracted_text = []
   
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text.append(shape.text)
   
    return "\n".join(extracted_text)

def load_docx(file):
    """
    Extracts text from a Word document (.docx).

    :param file_path: str, path to the Word document
    :return: str, extracted text
    """

    # Load the document
    document = Document(io.BytesIO(file.read()))
   
    # Extract text
    extracted_text = []
   
    for para in document.paragraphs:
        extracted_text.append(para.text)
   
    return "\n".join(extracted_text)

def load_txt(file):
    """
    Extracts text from a text file (.txt).

    :param file_path: str, path to the text file
    :return: str, extracted text
    """    
    stringio = StringIO(file.getvalue().decode("utf-8"))
    # To read file as string:
    extracted_text = stringio.read()
    return extracted_text


def load_document(file):
    
    if file.name.endswith('.pdf'):
        text = load_pdf_minder(file)
    elif file.name.endswith('.docx'):
        text = load_docx(file)
    elif file.name.endswith('.txt'):
        text = load_txt(file)
    elif file.name.endswith('.pptx'):
        text = load_ppt(file)
    return text

def printOutput(output):
    print(json.dumps(output,sort_keys=True, indent=3))



def extract_fields(names, descriptions, exs, llm, source_text, language="English"):


    #16 elt 0 -> 15

    example_1 = {
        names[0]:exs[0][0], 
        names[1]:exs[0][1],
        names[2]:exs[0][2],
        names[3]:exs[0][3],
        names[4]:exs[0][4],
        names[5]:exs[0][5],
        names[6]:exs[0][6],
        names[7]:exs[0][7],
        names[8]:exs[0][8],
        names[9]:exs[0][9],
        names[10]:exs[0][10],
        names[11]:exs[0][11],
        names[12]:exs[0][12],
        names[13]:exs[0][13],
        names[14]:exs[0][14],
        names[15]:exs[0][15],
    }
    example_2 = {
        names[0]:exs[1][0], 
        names[1]:exs[1][1],
        names[2]:exs[1][2],
        names[3]:exs[1][3],
        names[4]:exs[1][4],
        names[5]:exs[1][5],
        names[6]:exs[1][6],
        names[7]:exs[1][7],
        names[8]:exs[1][8],
        names[9]:exs[1][9],
        names[10]:exs[1][10],
        names[11]:exs[1][11],
        names[12]:exs[1][12],
        names[13]:exs[1][13],
        names[14]:exs[1][14],
        names[15]:exs[1][15],
    }
    example_3 = {
        names[0]:exs[2][0], 
        names[1]:exs[2][1],
        names[2]:exs[2][2],
        names[3]:exs[2][3],
        names[4]:exs[2][4],
        names[5]:exs[2][5],
        names[6]:exs[2][6],
        names[7]:exs[2][7],
        names[8]:exs[2][8],
        names[9]:exs[2][9],
        names[10]:exs[2][10],
        names[11]:exs[2][11],
        names[12]:exs[2][12],
        names[13]:exs[2][13],
        names[14]:exs[2][14],
        names[15]:exs[2][15],
    }
    
    

    schema = Object(
        id="contract",
        description="Information about a contract",
        attributes = [
            Text(
                id=names[0],
                description=descriptions[0]
            ),
            Text(
                id=names[1],
                description=descriptions[1]
            ),
            Text(
                id=names[2],
                description=descriptions[2]
            ),
            Text(
                id=names[3],
                description=descriptions[3]
            ),
            Text(
                id=names[4],
                description=descriptions[4]
            ),
            Text(
                id=names[5],
                description=descriptions[5]
            ),
            Text(
                id=names[6],
                description=descriptions[6]
            ),
            Text(
                id=names[7],
                description=descriptions[7]
            ),
            Text(
                id=names[8],
                description=descriptions[8]
            ),
            Text(
                id=names[9],
                description=descriptions[9]
            ),
            Text(
                id=names[10],
                description=descriptions[10]
            ),
            Text(
                id=names[11],
                description=descriptions[11]
            ),
            Text(
                id=names[12],
                description=descriptions[12]
            ),
            Text(
                id=names[13],
                description=descriptions[13]
            ),
            Text(
                id=names[14],
                description=descriptions[14]
            ),
            Text(
                id=names[15],
                description=descriptions[15]
            )
        ],
        examples=[
            (
                source_text,
                [
                    example_1,
                    #example_2,
                    #example_3,
                ],
            )
        ]
    )



    chain = create_extraction_chain(llm, schema)
    result = chain.invoke(source_text)
    return result["data"]
    #return example_1
