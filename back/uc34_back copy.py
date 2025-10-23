import pandas as pd
from kor.extraction import create_extraction_chain
from kor.nodes import Object, Text, Number
from PyPDF2 import PdfReader
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import io





def load_pdf(file):
    """
    Extracts text from a PDF document.

    :param file_path: str, path to the PDF file
    :return: str, extracted text
    """
    # Create a PDF reader object
    reader = PdfReader(file)
   
    # Extract text
    extracted_text = []
   
    for page in reader.pages:
        extracted_text.append(page.extract_text())
   
    return "\n".join(extracted_text)

def load_pdf_minder(file):
    
    text = extract_text(file)
    return text


def load_ppt(file):
    """
    Extracts text from a PowerPoint presentation.

    :param file_path: str, path to the PowerPoint file
    :return: str, extracted text
    """

    # Load the presentation
    presentation = Presentation(io.BytesIO(file.read()))
   
    # Extract text
    extracted_text = []
   
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text.append(shape.text)
   
    return "\n".join(extracted_text)

def load_docx(file):
    """
    Extracts text from a Word document (.docx).

    :param file_path: str, path to the Word document
    :return: str, extracted text
    """

    # Load the document
    document = Document(io.BytesIO(file.read()))
   
    # Extract text
    extracted_text = []
   
    for para in document.paragraphs:
        extracted_text.append(para.text)
   
    return "\n".join(extracted_text)

def load_txt(file):
    """
    Extracts text from a text file (.txt).

    :param file_path: str, path to the text file
    :return: str, extracted text
    """    
    stringio = StringIO(file.getvalue().decode("utf-8"))
    # To read file as string:
    extracted_text = stringio.read()
    return extracted_text


def load_document(file):
    
    if file.name.endswith('.pdf'):
        text = load_pdf_minder(file)
    elif file.name.endswith('.docx'):
        text = load_docx(file)
    elif file.name.endswith('.txt'):
        text = load_txt(file)
    elif file.name.endswith('.pptx'):
        text = load_ppt(file)
    return text

def printOutput(output):
    print(json.dumps(output,sort_keys=True, indent=3))



def extract_fields(names, descriptions, exs, llm, source_text, language="English"):


    attributes=[]
    example_1 = {}
    example_2 = {}
    example_3 = {}
    

    
    for i, elt in enumerate(names):
        example_1[elt] = exs[0][i]
        example_2[elt] = exs[1][i]
        example_3[elt] = exs[2][i]
    
    #examples = [example_1,example_2,example_3]

    schema = Object(
        id="contract",
        description="Information about a contract",
        attributes =  [Text(id=elt,description=descriptions[i]) for i, elt in enumerate(names)],
        examples=[
            (
                source_text,
                [
                    example_1,
                    example_2,
                    example_3,
                ],
            )
        ]
    )



    chain = create_extraction_chain(llm, schema)
    result = chain.invoke(source_text)
    return result["data"]
    #return schema
